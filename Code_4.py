# CODE 4

from pptx import Presentation
import os

def add_bullet_point_to_pptx(input_folder, output_folder, new_text_line):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if not filename.endswith(".pptx"):
            continue
       
        file_path = os.path.join(input_folder, filename)
        prs = Presentation(file_path)
        modified = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
               
                text_frame = shape.text_frame
                paragraphs = [p for p in text_frame.paragraphs if p.text.strip() != ""]
               
                if len(paragraphs) < 2:
                    # Skip if not a bullet list
                    continue
               
                last_paragraph = paragraphs[-1]

                # Determine font style from existing text
                font_name = None
                font_size = None
                if last_paragraph.runs:
                    font_name = last_paragraph.runs[0].font.name
                    font_size = last_paragraph.runs[0].font.size

                # Add empty paragraph to create a blank line gap
                blank_para = text_frame.add_paragraph()
                blank_para.text = " "

                # Add new paragraph maintaining formatting
                new_para = text_frame.add_paragraph()
                new_para.text = f"{len(paragraphs) + 1}. {new_text_line}"
                new_para.level = last_paragraph.level

                # Apply previous formatting
                if new_para.runs:
                    run = new_para.runs[0]
                    run.font.name = font_name if font_name else "Arial"
                    run.font.size = font_size

                modified = True

        # Save updates if any change was made
        if modified:
            output_path = os.path.join(output_folder, f"updated_{filename}")
            prs.save(output_path)
            print(f"✅ Updated file saved: {output_path}")
        else:
            print(f"⚠️ No bullet text found in: {filename}")

# Replace with your folders and text before running
input_folder = r"C:\Users\OUTPUT PPTS"
output_folder = r"C:\Users\Final"
new_text = "    SERVICE PART PER SPEC 49-00362-000 - NON-SERVICEABLE ASSEMBLY."

add_bullet_point_to_pptx(input_folder, output_folder, new_text)
