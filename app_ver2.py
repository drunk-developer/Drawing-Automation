import streamlit as st
import os
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from copy import deepcopy
import tempfile

# --- Helper Functions ---
REV_HEADERS = ["RELEASE NUMBER", "REV LTR", "REVISION DESCRIPTION", "BY", "DATE", "APPD"]

def is_revision_table(table):
    header = [cell.text.strip().upper() for cell in table.rows[0].cells]
    return header == REV_HEADERS

def find_balloons_and_texts_recursive(shapes):
    balloon_shapes = []
    text_shapes = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if sh.auto_shape_type in [9, 40, 56, 57]:
                    balloon_shapes.append(sh)
            except:
                pass
        if sh.has_text_frame and sh.text.strip():
            text_shapes.append(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            bs, ts = find_balloons_and_texts_recursive(sh.shapes)
            balloon_shapes.extend(bs)
            text_shapes.extend(ts)
    return balloon_shapes, text_shapes

def get_balloon_letters_flexible(slide):
    balloon_shapes, text_shapes = find_balloons_and_texts_recursive(slide.shapes)
    balloon_letters = []
    for balloon in balloon_shapes:
        bx, by, bw, bh = balloon.left, balloon.top, balloon.width, balloon.height
        b_center = (bx + bw / 2, by + bh / 2)
        nearest_letter, nearest_dist = "", float("inf")
        for txt in text_shapes:
            tx, ty, tw, th = txt.left, txt.top, txt.width, txt.height
            t_center = (tx + tw / 2, ty + th / 2)
            dist = ((b_center[0] - t_center[0]) ** 2 + (b_center[1] - t_center[1]) ** 2) ** 0.5
            text = txt.text.strip()
            if len(text) == 1 and dist < min(bw, bh):
                if dist < nearest_dist:
                    nearest_letter = text
                    nearest_dist = dist
        balloon_letters.append(nearest_letter)
    return balloon_letters

def extract_revision_data_multisheet_from_files(uploaded_files):
    per_drawing = {}
    with tempfile.TemporaryDirectory() as temp_dir:
        for uploaded_file in uploaded_files:
            ppt_path = os.path.join(temp_dir, uploaded_file.name)
            with open(ppt_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            prs = Presentation(ppt_path)
            drawing_name = os.path.splitext(uploaded_file.name)[0]
            revision_rows = []
            balloon_letters = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table and is_revision_table(shape.table):
                        table = shape.table
                        for i in range(1, len(table.rows)):
                            row_data = [cell.text.strip() for cell in table.rows[i].cells]
                            revision_rows.append(row_data)
                        balloon_letters = get_balloon_letters_flexible(slide)
                        while len(balloon_letters) < len(revision_rows):
                            balloon_letters.append("")
                        balloon_letters = balloon_letters[:len(revision_rows)]
            sheet_rows = []
            for row, balloon in zip(revision_rows, balloon_letters):
                sheet_rows.append(row + [balloon])
            if sheet_rows:
                columns = REV_HEADERS + ["Balloon Text"]
                per_drawing[drawing_name] = pd.DataFrame(sheet_rows, columns=columns)
    # Return Excel as a bytes object for online download
    with tempfile.NamedTemporaryFile(suffix=".xlsx") as tmp:
        with pd.ExcelWriter(tmp.name) as writer:
            for name, df in per_drawing.items():
                sheet_name = str(name)[:31]
                df.to_excel(writer, index=False, sheet_name=sheet_name)
        tmp.seek(0)
        return tmp.read() if per_drawing else None

TABLE_FONT = "Arial Narrow"
TABLE_SIZE = Pt(7)
BALLOON_FONT = "Arial"
BALLOON_SIZE = Pt(11)
REVISION_HEADERS = ["RELEASE NUMBER", "REV LTR", "REVISION DESCRIPTION", "BY", "DATE", "APPD"]

def update_table_and_balloon_for_all(excel_bytes, uploaded_ppt_files):
    results = {}
    with tempfile.NamedTemporaryFile(suffix=".xlsx") as xl_file, tempfile.TemporaryDirectory() as tmp_dir:
        xl_file.write(excel_bytes)
        xl_file.flush()
        xl = pd.ExcelFile(xl_file.name)
        for ppt_file in uploaded_ppt_files:
            ppt_name = os.path.splitext(ppt_file.name)[0]
            if ppt_name not in xl.sheet_names:
                continue
            df = xl.parse(ppt_name)
            df = df.dropna(how='all')
            if df.empty:
                continue
            revision_data = df[REVISION_HEADERS].values.tolist()
            balloon_values = df["Balloon Text"].dropna()
            balloon_letter = str(balloon_values.iloc[-1]) if not balloon_values.empty else ""
            pptx_path = os.path.join(tmp_dir, ppt_file.name)
            with open(pptx_path, "wb") as f:
                f.write(ppt_file.getbuffer())
            prs = Presentation(pptx_path)
            revision_done = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        headers = [cell.text.strip().upper() for cell in shape.table.rows[0].cells]
                        if is_revision_table(headers) and not revision_done:
                            while len(shape.table.rows) > 1:
                                shape.table._tbl.remove(shape.table.rows[1]._tr)
                            for rev in revision_data:
                                header_xml = shape.table.rows[0]._tr
                                new_xml = deepcopy(header_xml)
                                shape.table._tbl.append(new_xml)
                                new_row = shape.table.rows[len(shape.table.rows) - 1]
                                for i, val in enumerate(rev):
                                    if i < len(new_row.cells):
                                        cell = new_row.cells[i]
                                        cell.text = str(val) if val is not None else ""
                                        para = cell.text_frame.paragraphs[0]
                                        para.font.name = TABLE_FONT
                                        para.font.size = TABLE_SIZE
                            revision_done = True
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        txt = shape.text.strip()
                        if (len(txt) == 1 and txt.isalpha()) or (len(txt) == 2 and txt[0].isalpha() and txt[1] == '.'):
                            if balloon_letter:
                                shape.text = str(balloon_letter)
                                para = shape.text_frame.paragraphs[0]
                                para.font.name = BALLOON_FONT
                                para.font.size = BALLOON_SIZE
            # Save updated pptx as bytes
            out_bytes = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            prs.save(out_bytes.name)
            out_bytes.seek(0)
            results[ppt_file.name] = out_bytes.read()
            out_bytes.close()
    return results

def add_bullet_point_to_pptx(uploaded_ppt_files, new_text_line):
    results = {}
    for ppt_file in uploaded_ppt_files:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(ppt_file.getbuffer())
            tmp.flush()
            prs = Presentation(tmp.name)
            modified = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame
                    paragraphs = [p for p in text_frame.paragraphs if p.text.strip() != ""]
                    if len(paragraphs) < 2:
                        continue
                    last_paragraph = paragraphs[-1]
                    font_name = None
                    font_size = None
                    if last_paragraph.runs:
                        font_name = last_paragraph.runs[0].font.name
                        font_size = last_paragraph.runs[0].font.size
                    blank_para = text_frame.add_paragraph()
                    blank_para.text = " "
                    new_para = text_frame.add_paragraph()
                    new_para.text = f"{len(paragraphs) + 1}. {new_text_line}"
                    new_para.level = last_paragraph.level
                    if new_para.runs:
                        run = new_para.runs[0]
                        run.font.name = font_name if font_name else "Arial"
                        run.font.size = font_size
                    modified = True
            if modified:
                prs.save(tmp.name)
                tmp.seek(0)
                results[ppt_file.name] = tmp.read()
    return results

# ================= Streamlit UI =================

st.title("PowerPoint Engineering Drawings Automation (Cloud Compatible)")

st.info("Before using, please manually convert all .ppt files to .pptx format locally using Microsoft Office PowerPoint.")

stage = st.sidebar.radio(
    "Choose Automation Stage",
    (
        "Step 1: Extract Revision Data to Excel",
        "Step 2: Edit PPTX from Excel",
        "Step 3: Add Bullet Point to PPTX"
    )
)

if stage == "Step 1: Extract Revision Data to Excel":
    st.header("Step 1: Extract Revision Table and Balloon Data")
    uploaded_pptxs = st.file_uploader("Upload .pptx files", type="pptx", accept_multiple_files=True)
    if st.button("Extract Data to Excel"):
        if uploaded_pptxs:
            excel_bytes = extract_revision_data_multisheet_from_files(uploaded_pptxs)
            if excel_bytes:
                st.success("Extraction complete! Download your Excel file below.")
                st.download_button("Download Excel", data=excel_bytes, file_name="Extracted_Revision_Data.xlsx")
            else:
                st.warning("No revision or balloon data found in uploaded files.")

elif stage == "Step 2: Edit PPTX from Excel":
    st.header("Step 2: Edit PPTX Files Based on Excel Data")
    uploaded_excel = st.file_uploader("Upload updated Excel file", type="xlsx")
    uploaded_pptxs = st.file_uploader("Upload corresponding .pptx files for editing", type="pptx", accept_multiple_files=True)
    if st.button("Apply Edits to PPTX"):
        if uploaded_excel and uploaded_pptxs:
            excel_bytes = uploaded_excel.getbuffer()
            updated_files = update_table_and_balloon_for_all(excel_bytes, uploaded_pptxs)
            st.success(f"Updated {len(updated_files)} PPTX files! Download below.")
            for fname, file_bytes in updated_files.items():
                st.download_button(f"Download {fname}", data=file_bytes, file_name=f"updated_{fname}")

elif stage == "Step 3: Add Bullet Point to PPTX":
    st.header("Step 3: Add Bullet Point")
    new_text_line = st.text_area("Enter new bullet point text")
    uploaded_pptxs = st.file_uploader("Upload .pptx files for bullet addition", type="pptx", accept_multiple_files=True)
    if st.button("Add Bullet Point"):
        if new_text_line and uploaded_pptxs:
            updated_files = add_bullet_point_to_pptx(uploaded_pptxs, new_text_line)
            st.success(f"Added bullet to {len(updated_files)} PPTX files! Download below.")
            for fname, file_bytes in updated_files.items():
                st.download_button(f"Download {fname}", data=file_bytes, file_name=f"updated_{fname}")
