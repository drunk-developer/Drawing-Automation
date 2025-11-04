# CODE 3
#EDIT THE PPTXs FROM UPDATED EXCEL



import os
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy

TABLE_FONT = "Arial Narrow"
TABLE_SIZE = Pt(7)
BALLOON_FONT = "Arial"
BALLOON_SIZE = Pt(11)
REVISION_HEADERS = [
    "RELEASE NUMBER", "REV LTR", "REVISION DESCRIPTION", "BY", "DATE", "APPD"
]

def is_revision_table(headers):
    normalized = [h.upper().replace(" ", "") for h in headers]
    required = [h.upper().replace(" ", "") for h in REVISION_HEADERS]
    return all(h in normalized for h in required)

def clear_table_rows(table):
    while len(table.rows) > 1:
        table._tbl.remove(table.rows[1]._tr)

def add_revision_rows(table, revision_data):
    for rev in revision_data:
        header_xml = table.rows[0]._tr
        new_xml = deepcopy(header_xml)
        table._tbl.append(new_xml)
        new_row = table.rows[len(table.rows) - 1]
        for i, val in enumerate(rev):
            if i < len(new_row.cells):
                cell = new_row.cells[i]
                cell.text = str(val) if val is not None else ""
                para = cell.text_frame.paragraphs[0]
                para.font.name = TABLE_FONT
                para.font.size = TABLE_SIZE

def update_table_and_balloon_for_all(multisheet_excel, ppt_folder, output_folder):
    xl = pd.ExcelFile(multisheet_excel)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    for ppt_file in sorted(os.listdir(ppt_folder)):
        if ppt_file.lower().endswith(".pptx"):
            ppt_name = os.path.splitext(ppt_file)[0]
            try:
                if ppt_name not in xl.sheet_names:
                    print(f"Sheet for {ppt_file} not found, skipping.")
                    continue
                df = xl.parse(ppt_name)
                df = df.dropna(how='all')
                if df.empty:
                    print(f"Sheet {ppt_name} is empty, skipping.")
                    continue
                revision_data = df[REVISION_HEADERS].values.tolist()
                balloon_values = df["Balloon Text"].dropna()
                # Always cast to string
                balloon_letter = str(balloon_values.iloc[-1]) if not balloon_values.empty else ""
                pptx_path = os.path.join(ppt_folder, ppt_file)
                output_path = os.path.join(output_folder, ppt_file)
                prs = Presentation(pptx_path)
                revision_done = False
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_table:
                            headers = [cell.text.strip().upper() for cell in shape.table.rows[0].cells]
                            if is_revision_table(headers) and not revision_done:
                                clear_table_rows(shape.table)
                                add_revision_rows(shape.table, revision_data)
                                revision_done = True
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            txt = shape.text.strip()
                            if (len(txt) == 1 and txt.isalpha()) or (len(txt) == 2 and txt[0].isalpha() and txt[1] == '.'):
                                if balloon_letter:
                                    shape.text = str(balloon_letter)
                                    para = shape.text_frame.paragraphs[0]
                                    para.font.name = BALLOON_FONT
                                    para.font.size = BALLOON_SIZE
                prs.save(output_path)
                print(f"Updated: {ppt_file}")
            except Exception as e:
                print(f"Error updating {ppt_file}: {e}")

# Usage:
update_table_and_balloon_for_all(
    r"C:\Users\Extracted_excel.xlsx",
    r"C:\Users\INPUT PPTS",
    r"C:\Users\OUTPUT PPTS"
)
