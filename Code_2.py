# CODE 2
#EXTRACT REVISION TABLE AND BALLOON DATA

import os
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REV_HEADERS = [
    "RELEASE NUMBER", "REV LTR",
    "REVISION DESCRIPTION", "BY", "DATE", "APPD"
]

def is_revision_table(table):
    header = [cell.text.strip().upper() for cell in table.rows[0].cells]
    return header == REV_HEADERS

def find_balloons_and_texts_recursive(shapes):
    balloon_shapes = []
    text_shapes = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if sh.auto_shape_type in [9, 40, 56, 57]:
                    balloon_shapes.append(sh)
            except Exception:
                pass
        if sh.has_text_frame and sh.text.strip():
            text_shapes.append(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            bs, ts = find_balloons_and_texts_recursive(sh.shapes)
            balloon_shapes.extend(bs)
            text_shapes.extend(ts)
    return balloon_shapes, text_shapes

def get_balloon_letters_flexible(slide):
    balloon_shapes, text_shapes = find_balloons_and_texts_recursive(slide.shapes)
    balloon_letters = []
    for balloon in balloon_shapes:
        bx, by, bw, bh = balloon.left, balloon.top, balloon.width, balloon.height
        b_center = (bx + bw / 2, by + bh / 2)
        nearest_letter, nearest_dist = "", float("inf")
        for txt in text_shapes:
            tx, ty, tw, th = txt.left, txt.top, txt.width, txt.height
            t_center = (tx + tw / 2, ty + th / 2)
            dist = ((b_center[0] - t_center[0]) ** 2 + (b_center[1] - t_center[1]) ** 2) ** 0.5
            text = txt.text.strip()
            # Accept single character (letter or number)
            if len(text) == 1 and dist < min(bw, bh):
                if dist < nearest_dist:
                    nearest_letter = text
                    nearest_dist = dist
        balloon_letters.append(nearest_letter)
    return balloon_letters

def extract_revision_data_multisheet(input_folder, excel_path):
    # Dictionary to hold DataFrames per drawing
    per_drawing = {}
    for ppt_file in sorted(os.listdir(input_folder)):
        if ppt_file.lower().endswith(".pptx"):
            ppt_path = os.path.join(input_folder, ppt_file)
            prs = Presentation(ppt_path)
            drawing_name = os.path.splitext(ppt_file)[0]
            revision_rows = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table and is_revision_table(shape.table):
                        table = shape.table
                        for i in range(1, len(table.rows)):
                            row_data = [cell.text.strip() for cell in table.rows[i].cells]
                            revision_rows.append(row_data)
                        balloon_letters = get_balloon_letters_flexible(slide)
                        while len(balloon_letters) < len(revision_rows):
                            balloon_letters.append("")
                        balloon_letters = balloon_letters[:len(revision_rows)]
            sheet_rows = []
            for row, balloon in zip(revision_rows, balloon_letters):
                sheet_rows.append(row + [balloon])
            # Only create sheet if there is actual data
            if sheet_rows:
                columns = REV_HEADERS + ["Balloon Text"]
                per_drawing[drawing_name] = pd.DataFrame(sheet_rows, columns=columns)
    # Write all sheets to one Excel file
    with pd.ExcelWriter(excel_path) as writer:
        for name, df in per_drawing.items():
            # Sheet name max 31 chars
            sheet_name = str(name)[:31]
            df.to_excel(writer, index=False, sheet_name=sheet_name)
    print(f"Extraction complete. Sheets created per PPT in: {excel_path}")

# Usage
extract_revision_data_multisheet(
    r"C:\Users\INPUT PPTS",
    r"C:\Users\Extracted_excel.xlsx"
)